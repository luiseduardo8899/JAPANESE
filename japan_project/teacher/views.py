import os
import xml.etree.ElementTree as metree
from django.shortcuts import render
from django.http import HttpResponseRedirect, HttpResponse
from userinfo.models import Profile, Statistics
from userinfo.utils import *
from dictionary.utils import *
from dictionary.models import *
from languagebits.models import *
from languagebits.utils import *
from django.template import loader
from django.http import Http404
from django.urls import reverse
from django.core.files.storage import FileSystemStorage
from django.utils import timezone
from django.shortcuts import redirect 
import time
from csv import reader
import romkan
import datetime
import svgwrite
import random
from svgwrite import *
from cairosvg import svg2png
import logging #python logging utility
from django.conf import settings
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

pos_int_switch = {
"TBD": 0,
"ADJECTIVAL_NOUN": 1,
"ADJECTIVE": 2,
"ADJECTIVE_TARU": 3,
"ADVERB": 4,
"ADVERBIAL_NOUN": 5,
"CONJUNCTION": 6,
"COUNTER": 7, 
"EXPRESSION": 8, 
"FORMAL_NA_ADJECTIVE": 9, 
"INTERJECTION": 10, 
"NOUN": 11, 
"NOUN_VERB_ACTING_PRENOMINALLY": 12, 
"PREFIX": 13, 
"PRE_NOUN_ADJECTIVAL": 14, 
"PRONOUN": 15, 
"SPECIAL": 16, 
"SUFFIX": 17, 
"TAKES_NO": 18, 
"TAKES_SURU": 19, 
"TAKES_TO": 20, 
"TEMPORAl": 21, 
"USED_AS_PREFIX": 22, 
"USED_AS_SUFFIX": 23, 
"VERB5_U": 24, 
"VERB_1": 25, 
"VERB_1_ZURU": 26, 
"VERB_2_RU": 27, 
"VERB_5_BU": 28, 
"VERB_5_GU": 29, 
"VERB_5_KU": 30, 
"VERB_5_MU": 31, 
"VERB_5_RU": 32, 
"VERB_5_SU": 33, 
"VERB_5_TSU": 34, 
"VERB_ARCHAIC": 35, 
"VERB_AUXILIARY": 36, 
"VERB_INTRANSITIVE": 37, 
"VERB_IRREGULAR": 38, 
"VERB_RU": 39, 
"VERB_SPECIAL": 40, 
"VERB_SURU": 41, 
"VERB_TRANSITIVE": 42}

pos_text_switch = {
"TBD": {"TBD : UNDEFINED"} ,
"ADJECTIVAL_NOUN": {"adjectival nouns or quasi-adjectives (keiyodoshi)"} ,
"ADJECTIVE": {"adjective (keiyoushi)"} ,
"ADJECTIVE_TARU": {"`taru' adjective"} ,
"ADVERB": {"adverb (fukushi)"}  ,
"ADVERBIAL_NOUN": {"adverbial noun (fukushitekimeishi)"}  ,
"CONJUNCTION": {"conjunction"} ,
"COUNTER": {"counter"}  ,
"EXPRESSION": {"expressions (phrases, clauses, etc.)"} ,
"FORMAL_NA_ADJECTIVE": {"archaic/formal form of na-adjective"}  ,
"INTERJECTION": {"interjection (kandoushi)"}  ,
"NOUN": {"noun (common) (futsuumeishi)"}  ,
"NOUN_VERB_ACTING_PRENOMINALLY": {"noun or verb acting prenominally"}    ,
"PREFIX": {"prefix"} ,
"PRE_NOUN_ADJECTIVAL": {"pre-noun adjectival (rentaishi)"}  ,
"PRONOUN": {"pronoun"} ,
"SPECIAL": {"Special class of the verb used"}  ,
"SUFFIX": {"suffix"}  ,
"TAKES_NO": {"nouns which may take the genitive case particle `no'"}  ,
"TAKES_SURU": {"noun or participle which takes the aux. verb suru"}  ,
"TAKES_TO": {"adverb taking the `to' particle"}  ,
"TEMPORAl": {"noun (temporal) (jisoumeishi)"} ,
"USED_AS_PREFIX": {"noun, used as a prefix"}  ,
"USED_AS_SUFFIX": {"noun, used as a suffix"}  ,
"VERB5_U": {"Godan verb with `u' ending"}  ,
"VERB_1": {"Ichidan verb"}  ,
"VERB_1_ZURU": {"Ichidan verb - zuru verb (alternative form of -jiru verbs)"}  ,
"VERB_2_RU": {"Nidan verb (lower class) with `ru' ending (archaic)"}  ,
"VERB_5_BU": {"Godan verb with `bu' ending"} ,
"VERB_5_GU": {"Godan verb with `gu' ending"} ,
"VERB_5_KU": {"Godan verb with `ku' ending"} ,
"VERB_5_MU": {"Godan verb with `mu' ending"} ,
"VERB_5_RU": {"Godan verb with `ru' ending"} ,
"VERB_5_SU": {"Godan verb with `su' ending"}  ,
"VERB_5_TSU": {"Godan verb with `tsu' ending"}  ,
"VERB_ARCHAIC": {"Archaic form of the verb used"}  ,
"VERB_AUXILIARY": {"auxiliary verb"}  ,
"VERB_INTRANSITIVE": {"intransitive verb"}  ,
"VERB_IRREGULAR": {"irregular form of the verb conjugation used"} ,
"VERB_RU": {"irregular ru verb, plain form ends with -ri"} ,
"VERB_SPECIAL": {"Special class of the verb used"}  ,
"VERB_SURU": {"suru verb - special class"}  ,
"VERB_TRANSITIVE": {"transitive verb"} }


# Upload Vocabulary in XML format / Follow JMdict format


logging.basicConfig(filename='/home/luis/dictionaryViews.log', filemode='w', level=logging.DEBUG)
logger = logging.getLogger('/home/luis/dictionaryViewsLogger.log')
logger2 = logging.getLogger('/home/luis/dictionaryViewsLogger2.log')
#namesapce required for xml:lang in lsource
nsmap = {"xml": "http://www.w3.org/XML/1998/namespace"}

def create_videos(request):
    error_in_file = False #TODO: Check PPTX creation worked
    return render(request, 'teacher/create_videos.html', {
        'error_in_file': error_in_file
    })

def create_powerpoint(request):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    COLOR1 = RGBColor(0x2B, 0xD6, 0xAB)
    COLOR2 = RGBColor(0xE6, 0xD7, 0x2E)
    COLOR3 = RGBColor(0xF5, 0x4E, 0x30)
    BLACK1 = RGBColor(0x00, 0x00, 0x00)
    WHITE1 = RGBColor(0xFF, 0xFF, 0xFF)
    
    title.text = "GOKOKAN"
    subtitle.text = "Online Language Learning"

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = "Online Language Learning"

    #Add a shape
    shapes = slide.shapes
    left = top = width = height = Inches(1.0)
    shape = shapes.add_shape( MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLOR1
    line = shape.line
    line.fill.background() #make it invisible
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "A"
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(60)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = BLACK1


    width = height = Inches(1.0)
    left = top = Inches(2.0)
    shape = shapes.add_shape( MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLOR2
    line = shape.line
    line.fill.background() #make it invisible
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "B"
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(60)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = WHITE1

    #TODO: Continuar AQUI -> Convertir en una funcion!
    #TODO: Medir el texto? Ratio text size to shape width/height?
    height = Inches(1.0)
    width = Inches(8.0)
    left = Inches(1.0)
    top = Inches(4.0)
    shape = shapes.add_shape( MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLOR3
    line = shape.line
    line.fill.background() #make it invisible
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "ありがとうございます！"
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(50)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = WHITE1

    prs.save('static/hello.pptx')

    error_in_file = False #TODO: Check PPTX creation worked
    return render(request, 'teacher/successful_upload.html', {
        'uploaded_file_url': "static/hello.pptx",
        'error_in_file': error_in_file
    })


def create_images(request):
    user = get_user(request)
    template = loader.get_template('teacher/create_images.html')

    #loop through all images
    #random number from 1 -100000
    for entry in  Entry.objects.all():
        #pkid = random.randint(0, Entry.objects.count() - 1)
        # pkid = random.randint(1,100000)

        #entry = get_entry_by_id(pkid)
        pkid = entry.pk
        text = entry.get_text()

        bk = random.randint(1,2)
        #furigana = entry.get_furigana()
        romanji = entry.get_romanji()
        definition = entry.get_definition()
        if pkid % 3 == 0 :
            bkcolor = "#2bd6ab" 
        elif pkid % 2 == 0 :
            bkcolor = "#e6d72e" 
        else :
            bkcolor = "#f54e30"

        if bk == 1:
            color1 = "white"
            color2 = "black"
        else: 
            color1 = "black"
            color2 = "white"

        #TODO: Check if user is admin
        #dwg = svgwrite.Drawing('static/test2.svg', height=900, width=900)
        imgname = "static/img_sm_"+str(pkid)+".svg"
        #dwg = svgwrite.Drawing(filename = "static/test2.svg", size=("550px", "270px"))
        #dwg.add_stylesheet('static/svgstyle.css', title="svgstyle") # same rules as for html files
        #dwg.add(dwg.rect((0, 0), (550, 135), fill=bkcolor))
        #g = dwg.g(class_="myclass")
        #g.add(dwg.text(text, insert=(40,60), style = "font-size:40px; font-family:Arial", fill=color1))
        #g.add(dwg.text("[ "+romanji+" ]", insert=(40,100), style = "font-size:20px; font-family:Arial;", fill=color1))

        #lendef = len(definition)
        #if lendef > 50:
        #    definition = definition[0:50]+'...'
        #g.add(dwg.text(definition, insert=(30,160), style = "font-size:20px; font-family:Arial", fill="black"))

        #g.add(dwg.text("gokokan.com", insert=(400,250), style = "font-size:15px; font-family:Arial", fill="black"))

        #dwg.add(g)
        #dwg.save()

        #svg2png(url="static/test2.svg", write_to="static/myfile.png")


        #--------------------------------- Facebook drawing ---------------------------------------------------------------#
        dwg2 = svgwrite.Drawing(filename = imgname, size=("260px", "260px"))
        dwg2.add_stylesheet('static/svgstyle.css', title="svgstyle") # same rules as for html files
        dwg2.add(dwg2.rect((0, 0), (260, 260), fill=bkcolor))
        g2 = dwg2.g(class_="myclass")
        g2.add(dwg2.text(text[0], insert=(80,140), style = "font-size:100px; font-family:Arial", fill=color1))

        g2.add(dwg2.text("GOKOKAN", insert=(60,220), style = "font-size:30px; font-family:Arial", fill=color2))

        dwg2.add(g2)
        dwg2.save()
        #--------------------------------- Facebook drawing ---------------------------------------------------------------#

    context = {
        'user': user,
    }
    return HttpResponse(template.render(context, request))

#Update Kanji List from /media/GK_grammar files:
def upload_kanji(request):
    user = get_user(request)
    if user.is_superuser :
        template = loader.get_template('teacher/upload_kanji.html')
        context = {
            'user': user,
        }
    else:
        template = loader.get_template('teacher/sorry.html')
        context = {
            'user': user,
        }
    return HttpResponse(template.render(context, request))

#process the XML file from /media/ directory and update the kanji in the database
def update_kanji(request):
    user = get_user(request)

    if user.is_superuser :
        file_ = open(os.path.join(settings.DATA_ROOT, 'GKkanji.xml'))

        error_in_file = False #TODO: Check XML format matches GKgrammar
        if error_in_file == False :
            process_kanji_xml_file(file_)

        uploaded_file_url = os.path.join(settings.DATA_ROOT, 'GKkanji.xml')
        return render(request, 'teacher/successful_upload.html', {
            'uploaded_file_url': uploaded_file_url,
            'error_in_file': error_in_file
        })
    else:
        return render(request, 'teacher/sorry.html', {
            'message': "Please head back to home page" #TODO: redirect to homepage?
        })

#Upload GKgrammar XML
def upload_grammar(request):
    user = get_user(request)
    profile = None
    #profile = get_profile(user)
    template = loader.get_template('teacher/upload_grammar.html')
    #TODO: Check if user is admin
    context = {
        'user': user,
        'profile': profile,
    }
    return HttpResponse(template.render(context, request))


#process form and update the vocabulary in the database
def update_grammar(request):
    user = get_user(request)
    #TODO: check that user is admin user
	#try and retrieve Form inputs..
    if request.method == 'POST' and request.FILES['xml_data_file']:
        xml_file = request.FILES['xml_data_file']
        fs = FileSystemStorage()
        filename = fs.save(xml_file.name, xml_file)
        uploaded_file_url = fs.url(filename)
        #process the csv data
        #python3 supports utf-8 encoding natively
        # error_in_file = check_csv_format(uploaded_file_url)
        #error_in_file = check_xml_format()
        error_in_file = False #TODO: Check XML format matches GKgrammar
        if error_in_file == False :
            #BASE_DIR = "/home/luis/TRIBU/projects/JAPANESE/"
            #MEDIA_URL = os.path.join(BASE_DIR, uploaded_file_url)
            process_grammar_xml_file(uploaded_file_url)
            #process_csv_file(uploaded_file_url)

        return render(request, 'teacher/successful_upload.html', {
            'uploaded_file_url': uploaded_file_url,
            'error_in_file': error_in_file
        })


        # Always return an HttpResponseRedirect after successfully dealing
        # with POST data. This prevents data from being posted twice if a
        # user hits the Back button.
        # TODO: how to pass arguments to redirect page ? 
    return HttpResponseRedirect('teacher/successful_upload')


def upload_dictionary(request):
    user = get_user(request)
    profile = get_profile(user)
    template = loader.get_template('teacher/upload_vocab.html')
    #TODO: Check if user is admin
    context = {
        'user': user,
        'profile': profile,
    }
    return HttpResponse(template.render(context, request))

#process form and update the vocabulary in the database
def update_dictionary(request):
    user = get_user(request)
    #TODO: check that user is admin user
	#try and retrieve Form inputs..
    if request.method == 'POST' and request.FILES['xml_data_file']:
        xml_file = request.FILES['xml_data_file']
        fs = FileSystemStorage()
        filename = fs.save(xml_file.name, xml_file)
        uploaded_file_url = fs.url(filename)
        #process the csv data
        #python3 supports utf-8 encoding natively
        # error_in_file = check_csv_format(uploaded_file_url)
        #error_in_file = check_xml_format()
        error_in_file = False #TODO: Check XML format matches JMDICT
        if error_in_file == False :
            #BASE_DIR = "/home/luis/TRIBU/projects/JAPANESE/"
            #MEDIA_URL = os.path.join(BASE_DIR, uploaded_file_url)
            process_xml_file(uploaded_file_url)
            #process_csv_file(uploaded_file_url)

        return render(request, 'teacher/successful_upload.html', {
            'uploaded_file_url': uploaded_file_url,
            'error_in_file': error_in_file
        })


        # Always return an HttpResponseRedirect after successfully dealing
        # with POST data. This prevents data from being posted twice if a
        # user hits the Back button.
        # TODO: how to pass arguments to redirect page ? 
    return HttpResponseRedirect('teacher/successful_upload')

def upload_vocab(request):
    user = get_user(request)
    if user.is_anonymous : 
        return redirect("/account/login")
    profile = get_profile(user)
    template = loader.get_template('teacher/upload_vocab.html')
    #TODO: Check if user is admin
    context = {
        'user': user,
        'profile': profile,
    }
    return HttpResponse(template.render(context, request))

#process form and update the vocabulary in the database
def update_vocab(request):
    user = get_user(request)
    #TODO: check that user is admin user
	#try and retrieve Form inputs..
    if request.method == 'POST' and request.FILES['xml_data_file']:
        xml_file = request.FILES['xml_data_file']
        fs = FileSystemStorage()
        filename = fs.save(xml_file.name, xml_file)
        uploaded_file_url = fs.url(filename)
        #process the csv data
        #python3 supports utf-8 encoding natively
        # error_in_file = check_csv_format(uploaded_file_url)
        #error_in_file = check_xml_format()
        error_in_file = False #TODO: Check XML format matches JMDICT
        if error_in_file == False :
            #BASE_DIR = "/home/luis/TRIBU/projects/JAPANESE/"
            #MEDIA_URL = os.path.join(BASE_DIR, uploaded_file_url)
            process_vocab_xml_file(uploaded_file_url)
            #process_csv_file(uploaded_file_url)

        return render(request, 'teacher/successful_upload.html', {
            'uploaded_file_url': uploaded_file_url,
            'error_in_file': error_in_file
        })


        # Always return an HttpResponseRedirect after successfully dealing
        # with POST data. This prevents data from being posted twice if a
        # user hits the Back button.
        # TODO: how to pass arguments to redirect page ? 
    return HttpResponseRedirect('teacher/successful_upload')



def successful_upload(request):
    user = get_user(request)
    profile = get_profile(user)
    template = loader.get_template('teacher/successful_upload.html')
    #TODO: Check if user is admin
    context = {
        'user': user,
        'profile': profile,
    }
    return HttpResponse(template.render(context, request))

def check_csv_format(csv_file):
    error_in_file = False 
    with open(csv_file, encoding='utf-8') as f:
        counter = 0
        myreader = reader(f, delimiter=';')
        for row in myreader:
            counter+=1
            if counter == 1 and row[0] != 'kanji' : 
                logger.info('FORMAT ERROR: expecting \'kanji\' got: %s ' % row[0])
                error_in_file = True
            if counter == 1 and row[1] != 'furigana' : 
                logger.info('FORMAT ERROR: expecting \'furigana\' got: %s ' % row[1])
                error_in_file = True 
            if counter == 1 and row[2] != 'definition' : 
                logger.info('FORMAT ERROR: expecting \'definition\' got: %s ' % row[2])
                error_in_file = True 
            if counter < 10 :
                if(row[0] != ''):
                    logger.info('Kanji: %s' % row[0])
                if(row[1] != ''):
                    logger.info('Furigana: %s ' %row[1])
                else:
                    logger.info('CSV_ERROR: NO FURIGANA')
                    error_in_file = True 
                if(row[2] != ''):
                    logger.info('Definition: '+row[2]+'\n')
                else:
                    logger.info('CSV_ERROR: NO DEFINITION')
                    error_in_file = True 
    return error_in_file

#Add all Keb Elements found in the entry
def add_keb(entry, k_ele_set):
    for k_ele in k_ele_set :
        keb_set_x = k_ele.findall('keb')
        for keb_x in keb_set_x :
            #check if keb entry exists, and get a pointer or create a new one
            kebs = Keb.objects.all().filter(text = keb_x.text) #for entry to be duplicate the text must be exactly the same
            if not kebs :
                logger.info("IF DETECTED NO KEBS CREATE NEW ONE %s" % kebs)
                keb = Keb()
                keb.text = keb_x.text
                keb.save()
                logger.info('ADDED NEW KEB: %s' % keb.text)
            else :
                keb = kebs[0]
                logger.info('USING EXISTING KEB: %s ' % keb.text) 
            if len(kebs) > 1 : 
                logger.warning('WARNING DUPLICATE KEBS: %s ' % keb.text) 
            keb.entry.add(entry) # add entry pointer to list
            keb.save()

#Add all Reb Elements found in the entry
def add_reb(entry, r_ele_set):
    for r_ele in r_ele_set :
        reb_set_x = r_ele.findall('reb')
        for reb_x in reb_set_x :
            #check if reb entry exists, and get a pointer
            #.. for entry to be duplicate the text must be exactly the same
            rebs = Reb.objects.all().filter(furigana = reb_x.text) 
            if not rebs : #create a new one and save
                logger.info("IF DETECTED NO REBS CREATE NEW ONE %s " % rebs)
                reb = Reb()
                reb.furigana = reb_x.text
                roma = romkan.to_roma(reb_x.text) # automatic convert to romanji
                reb.romanji = roma
                reb.save()
                logger.info('ROMA CONVERSION: %s' % reb.romanji)
                logger.info('ADDED NEW REB: %s' % reb.furigana)
            else :
                reb = rebs[0]
            if len(rebs) > 1 : 
                logger.warning('WARNING DUPLICATE REBS: %s ' % reb.furigana) 
            reb.entry.add(entry) # add entry pointer to list
            reb.save()
            logger.info('FINISHED UPDATING REB: %s' % reb.furigana)

#TODO: Post processing , check for duplicate entries, double check the entry Keb and Reb, and Sense match, or similar

#Add all tag elements related to the input Sense
def add_sense_tag(sense, sense_x):
    #pos, misc, field, dial are all a single string based thus they all use SenseTag Model
    pos_set_x = sense_x.findall('pos')
    add_sense_tag_x(sense, pos_set_x, POS_TAG)

    misc_set_x = sense_x.findall('misc')
    add_sense_tag_x(sense, misc_set_x, MISC_TAG)

    field_set_x = sense_x.findall('field')
    add_sense_tag_x(sense, field_set_x, FIELD_TAG)

    dial_set_x = sense_x.findall('dial')
    add_sense_tag_x(sense, dial_set_x, DIALECT_TAG)

    #lsource and sinf tags have their own model
    lsource_set_x = sense_x.findall('lsource')
    add_lsource_x(sense, lsource_set_x)

    sinf_set_x = sense_x.findall('s_inf')
    add_sinf_x(sense, sinf_set_x)

def add_sinf_x(sense, sinf_set_x ):
    for sinf_x in sinf_set_x :
        #check if sense tag entry exists, and get a pointer 
        #.. for entry to be duplicate the text must be exactly the same
        logger.info('Task: add_sinf_x() : Called ... ')
        if sinf_x.text == "" :
            logger.warning('WARNING: Task: add_sinf_x() : Empty text in sinf_x.text: %s ' % sinf_x.text) 
        sinfs = SenseInf.objects.all().filter(text = sinf_x.text) 
        if not sinfs : #create a new one and save
            logger.info('Creating new : SenseInf for : %s' % sinf_x.text)
            senseinf = SenseInf()
            senseinf.text = sinf_x.text
            senseinf.save()
            logger.info('Task: add_sinf_x() : Adding Sense to SenseInf  : %s' % sinf_x.text)
            senseinf.sense.add(sense)
            senseinf.save()
        else: #default to first entry as any other entry would be duplicate
            logger.info('Task: add_sinf_x() : Adding Sense SenseInf  : %s' % sinf_x.text)
            senseinf = sinfs[0]
            senseinf.sense.add(sense)
            senseinf.save()
        if len(sinfs) > 1 : 
            logger.warning('WARNING: Task: add_sinf_x() : DUPLICATE SenseInf OBJECTS: %s ' % sinf_x.text) 

def add_lsource_x(sense, lsource_set_x):
    for lsource_x in lsource_set_x :
        #check if sense tag entry exists, and get a pointer 
        #.. for entry to be duplicate the text must be exactly the same
        logger.info('Task: add_lsource_x() : Called ... ')

        #text contains the vocabulary entry corresponding txt in the source language
        lsources = LSource.objects.all().filter(text = lsource_x.text) 
        if not lsources : #create a new one and save
            logger.info('Creating new : LSource for : %s' % lsource_x.text)
            lsource = LSource()
            if lsource_x.text == None:
                lsource.text =  "" # TODO: models: how can we allow a string to be empty ?
                logger.warning('lsource_x.text is empty %s' % lsource_x.text)
            else: 
                lsource.text = lsource_x.text

            if '{http://www.w3.org/XML/1998/namespace}lang' in lsource_x.attrib:
                lang = lsource_x.attrib['{http://www.w3.org/XML/1998/namespace}lang'] #source language
                logger.info('Detected langs %s ' % lang)
                lsource.lang = lang
            else:
                logger.info('NO {http://www.w3.org/XML/1998/namespace}lang detected')

            if 'ls_wasei' in lsource_x.attrib:
                lswasei = lsource_x.attrib['ls_wasei']
                logger.info('Detected lswasei %s ' % lswasei)
                lsource.lswasei = True

            if 'ls_type' in lsource_x.attrib:
                ls_type = lsource_x.attrib['ls_type']
                logger.info('Detected ls_type %s ' % ls_type)
                lsource.full = False
            lsource.save()
            logger.info('Task: add_lsource_x() : Adding Sense to LSource  : %s' % lsource_x.text)
            lsource.sense.add(sense)
            lsource.save()
        else: #default to first entry as any other entry would be duplicate
            logger.info('Task: add_lsource_x() : Adding Sense LSource  : %s' % lsource_x.text)
            lsource = lsources[0]
            lsource.sense.add(sense)
            lsource.save()
        if len(lsources) > 1 : 
            logger.warning('WARNING: Task: add_lsource_x() : DUPLICATE SenseInf OBJECTS: %s ' % lsource_x.text) 

def add_sense_tag_x(sense, tag_x_set, tag_type):
    for tag_x in tag_x_set :
        #check if sense tag entry exists, and get a pointer 
        #.. for entry to be duplicate the text must be exactly the same
        logger.info('Task: add_sense_tag_x() : Called with tag type : %s'  % tag_type)
        sensetags = SenseTag.objects.all().filter(tid = tag_type).filter(text = tag_x.text) 
        if not sensetags : #create a new one and save
            logger.info('Creating new : SenseTag for : %s' % tag_x.text)
            sensetag = SenseTag(tid=MISC_TAG)
            if tag_type == POS_TAG : 
                sensetag.tid = POS_TAG
            elif tag_type == MISC_TAG : 
                sensetag.tid = MISC_TAG
            elif tag_type == FIELD_TAG : 
                sensetag.tid = FIELD_TAG
            elif tag_type == DIALECT_TAG : 
                sensetag.tid = DIALECT_TAG
            else :
                logger.warning('Task: add_sense_tag_x() : Called with wrong tag type : %s'  % tag_type)
                sensetag.tid = MISC_TAG
            sensetag.text = tag_x.text
            sensetag.save()
            logger.info('Task: add_sense_tag_x() : Adding Sense to SenseTag  : %s' % tag_x.text)
            sensetag.sense.add(sense)
            sensetag.save()
        else: #default to first entry as any other entry would be duplicate
            logger.info('Task: add_sense_tag_x() : Adding Sense SenseTag  : %s' % tag_x.text)
            sensetag = sensetags[0]
            sensetag.sense.add(sense)
            sensetag.save()
        if len(sensetags) > 1 : 
            logger.warning('WARNING: Task: add_sense_tag_x() : DUPLICATE SenseTag OBJECTS: %s ' % tag_x.text) 


#Add all Sense Elements found in the entry
def add_sense(entry, sense_x_set):
    for sense_x in sense_x_set :
        gloss_set_x = sense_x.findall('gloss')
        sense = Sense()
        sense.save()
        sense.entry.add(entry)
        sense.save()
        
        add_sense_tag(sense, sense_x)

        for gloss_x in gloss_set_x :
            glosses = Gloss.objects.all().filter(text = gloss_x.text) 
            if not glosses : #create a new one and save
                gloss = Gloss()
                #slice gloss elemetn if too long
                if len(gloss_x.text) > 199:
                    gloss.text = gloss_x.text[0:199]
                else:
                    gloss.text = gloss_x.text
                logger.warning('WARNING: Had to slice a Gloss Element: %s ' % gloss.text) 
                gloss.save()
            else :
                gloss = glosses[0]
            if len(glosses) > 1 : 
                logger.warning('WARNING DUPLICATE GLOSS OBJECTS: %s ' % gloss.text) 
            gloss.sense.add(sense) # add entry pointer to list
            gloss.save()
            logger.info('FINISHED UPDATING GLOSS: %s' % gloss.text)
           
#TODO: Save a backup XML of the GKgrammar in the database
def dump_grammar_xml_file():
    return "TODO"


#Add all Grammar Description Elements found in the GrammarEntry
def add_grammar_description(entry, description_set):
    for description_x in description_set :
        gdescription = GrammarDescription()
        gdescription.text = description_x.text
        gdescription.save()
        gdescription.entry.add(entry) # add entry pointer to list
    

def symbol_get_description(symbol):
    return "TBD-Symbol-Description"

#Add all GrammarPattern Elements found in the GrammarEntry
def add_grammar_formula(entry, formula_set):
    for formula_x in formula_set:
        name_x = formula_x.findall('fname')
        pformulas = PatternFormula.objects.all().filter(text = name_x[0].text) #for entry to be duplicate the text must be exactly the same
        if not pformulas :
            logger.info("No duplicates detected, Creating a new PatternFomrula")
            pformula = PatternFormula()
            pformula.text = name_x[0].text
            pformula.save()
            logger.info('ADDED NEW PatternFormula: %s' % pformula.text)
            item_set_x = formula_x.findall('pitem') #Multiple PatternItems form a PatternFormula
            for item_x in item_set_x :
                pitems = PatternItem.objects.all().filter(text = item_x.text) #for entry to be duplicate the text must be exactly the same
                if not pitems :
                    logger.info("IF DETECTED NO PITEM CREATE NEW ONE" )
                    pitem = PatternItem()
                    pitem.symbol = item_x.text
                    pitem.text = symbol_get_description(item_x.text)
                    pitem.save()
                    logger.info('ADDED NEW PatternItem: %s' % pitem.text)
                else :
                    pitem = pitems[0]
                    logger.info('USING EXISTING PatternItem: %s ' % pitem.text) 
                if len(pitems) > 1 : 
                    logger.warning('WARNING DUPLICATE PatternItems: %s ' % pitem.text) 
                pitem.formula.add(pformula) # add entry pointer to list
                pitem.save()
        else :
            pformula = pformulas[0]
            logger.info('WARNING:  Detected a Duplicate GrammarFormula in a Grammar Entry: %s. Ignoring this FormulaEntry' % pformula.text)

        pformula.entry.add(entry) # add entry pointer to list

#Process a GKgrammar file
def process_grammar_xml_file(xml_file):
    list_of_entries = []
    number = 0

    logger2.info('Processing XML File - ')
    #Main execution code
    #load the entire GKgrammar XML File 
    tree = metree.parse(xml_file)

    #get root of the XML File
    root = tree.getroot()

    #TODO: Check the XML format...
    logger2.info('TODO: Check XML format  - ')

    #get all entries in the file
    entries_x = root.findall('entry')
    logger.info("Got root %s" % root)
    for entry_x in entries_x:
        number += 1
        #create the entry
        level_x = 0 #default level 0, level assigned in later processing
        entry = GrammarEntry()
        ent_seq_set_x = entry_x.findall('ent_seq')
        name_set_x = entry_x.findall('name')
        entry.seqid = ent_seq_set_x[0].text
        entry.text = name_set_x[0].text
        #entry.pub_date = datetime.datetime.now()
        #get kanji and reading element, kanji element may be empty
        jlpt_level_set = entry_x.findall('jlpt') # Kanji Element
        entry.jlpt = jlpt_level_set[0].text
        level_set = entry_x.findall('level') # Sublevel
        entry.level = jlpt_level_set[0].text
        summary_set = entry_x.findall('summary') # Summary
        entry.summary = summary_set[0].text
        entry.save()

        description_set = entry_x.findall('description') # Full Description
        formula_set = entry_x.findall('formula') # Grammar Pattern

        add_grammar_description(entry, description_set)
        add_grammar_formula(entry, formula_set)
        
    return list_of_entries # return the list of uploaded entries ( KEB/REB )


#Format for Vocab XML Files is defined in GKVocab 
def process_vocab_xml_file(xml_file):
    list_of_entries = []
    number = 0
    counter = 0
    level_counter = 1

    logger2.info('Processing GKVocab XML File - ')
    #Main execution code
    #load the entire JMdict_e XML File 
    tree = metree.parse(xml_file)

    #get root of the XML File
    root = tree.getroot()

    #TODO: Check the XML format... ( Check fileds are all properly defined )
    logger2.info('TODO: Check XML format  - ')

    #get all entries in the file
    entries_x = root.findall('vocab_entry')
    logger.info("Got root %s" % root)
    for entry_x in entries_x:
        number += 1
        #create the entry
        level_x = 0 #default level 0, level assigned in later processing

        #Check if entry already exists
        #Else
        #get kanji and reading element, kanji element may be empty
        kanji_set = entry_x.findall('kanji') # Kanji Element
        furigana_set = entry_x.findall('furigana') # Reading Element ( furigana )
        ent_seq_set_x = entry_x.findall('ent_seq')
        dict_seq_set_x = entry_x.findall('dict_seq')
        misc_set_x = entry_x.findall('misc')
        jlpt_set = entry_x.findall('jlpt') # Kanji Element
        level_set = entry_x.findall('level') # Reading Element ( furigana )
        definition_set = entry_x.findall('definition') # Sense element
        pos_set = entry_x.findall('pos') # Sense element

        vocabs = Vocabulary.objects.all().filter(seqid = ent_seq_set_x[0].text)

        if len(vocabs) > 0 and vocabs != None :
            vocab = vocabs[0]
        else :
            vocab = Vocabulary()
            vocab.seqid = ent_seq_set_x[0].text
            vocab.pub_date = datetime.datetime.now()
            if kanji_set[0].text == "" or kanji_set[0].text == None:
                vocab.text = ""
                vocab.furigana = furigana_set[0].text
            else:
                vocab.text = kanji_set[0].text
                vocab.furigana = furigana_set[0].text
            vocab.save()

        vocab = vocabs[0]

        #Update Kanji and Furigana
        if kanji_set[0].text == "" or kanji_set[0].text == None:
            vocab.text = ""
            vocab.furigana = furigana_set[0].text
        else:
            vocab.text = kanji_set[0].text
            vocab.furigana = furigana_set[0].text


        #Get JLPT Level and overall level of difficulty
        vocab.jlpt = jlpt_set[0].text

        #If no level set, assign one arbitrarily
        if not level_set or level_set[0].text == "" or level_set[0].text == None:
            #TODO: For now increase level based on 5 entries per level
            counter += 1
            if counter == 5:
                level_counter += 1
                counter = 0
            vocab.level = level_counter
        else:
            vocab.level = level_set[0].text

        #Add Romanji 
        roma = romkan.to_roma(furigana_set[0].text) # automatic convert to romanji
        vocab.romanji = roma

        #get sense elements, these include tags and definitions
        #TODO: LUIS ENABLE definition add_definition(vocab, definition_set)
        #add_definition(vocab, definition_set)
        add_pos(vocab, pos_set)
        #TODO: Add MISC / S_INF

        misc_set = entry_x.findall('misc') # Sense element
        for misc_x in misc_set :
            if misc_x.text == "KANA_ALONE":
                vocab.usu_kana = True

        vocab.save()
        
    return list_of_entries # return the list of uploaded entries ( KEB/REB )

#Add all Definition Elements found in the Vocabulary Entry
def add_pos(entry, pos_set):
    for pos_x in pos_set :
        #check if Definition entry exists, and get a pointer
        #.. for entry to be duplicate the text must be exactly the same
        poss = PartOfSpeech.objects.all().filter(text = pos_x.text) 
        if not poss : #create a new one and save
            logger.info("IF DETECTED NO PRE EXISTING PartOfSpeech Objects %s " % poss)
            pos1 = PartOfSpeech()
            pos1.pos = pos_int_switch.get(pos_x.text)
            pos1.text = pos_text_switch.get(pos_x.text)
            pos1.save()
            logger.info('ADDED NEW PARTOFSPEECH: %s %s' % (pos1.pos, pos1.text) )
        else :
            pos1 = poss[0]
        if len(poss) > 1 : 
            logger.warning('WARNING DUPLICATE PARTOFSPEECH OBJECTS: %s %s' % (pos1.pos, pos1.text)) 

        #Make sure the POS entry does not already exist in the entry
        entry_pos_set = entry.part_of_speech.all()
        pos_already_set = 0
        if len(entry_pos_set) > 0:
            for entry_pos_x in entry_pos_set:
                if entry_pos_x.pos == pos1.pos:
                    pos_already_set = 1

        if pos_already_set == 0:
            entry.part_of_speech.add(pos1) # add entry pointer to list
            entry.save()
            pos1.save()
        logger.info('FINISHED UPDATING PARTOFSPEECH: %s' % pos1.text)

#Add all Definition Elements found in the Vocabulary Entry
def add_definition(entry, def_set):
    for def_x in def_set :
        #check if Definition entry exists, and get a pointer
        #.. for entry to be duplicate the text must be exactly the same
        defs = LangDefinition.objects.all().filter(entext = def_x.text) 
        if not defs : #create a new one and save
            logger.info("IF DETECTED NO PRE EXISTING LANGDEFINITIONS CREATE NEW ONE %s " % defs)
            def1 = LangDefinition()
            def1.entext = def_x.text
            def1.save()
            logger.info('ADDED NEW LANGDEFINITION: %s' % def1.entext)
        else :
            def1 = defs[0]
        if len(defs) > 1 : 
            logger.warning('WARNING DUPLICATE LANGDEFINITIONS: %s ' % def1.entext) 
        def1.entry.add(entry) # add entry pointer to list
        def1.save()
        logger.info('FINISHED UPDATING LANGUAGEDEFINITION: %s' % def1.entext)

#Format for XML Files will be same as used in the JMdct file. 
#Any nes aditions to the vocabulary should be done manually or using the same format as JMdict
def process_xml_file(xml_file):
    list_of_entries = []
    number = 0

    logger2.info('Processing XML File - ')
    #Main execution code
    #load the entire JMdict_e XML File 
    tree = metree.parse(xml_file)

    #get root of the XML File
    root = tree.getroot()

    #TODO: Check the XML format...
    logger2.info('TODO: Check XML format  - ')

    #get all entries in the file
    entries_x = root.findall('entry')
    logger.info("Got root %s" % root)
    for entry_x in entries_x:
        number += 1
        #create the entry
        level_x = 0 #default level 0, level assigned in later processing
        entry = VocabEntry()
        ent_seq_set_x = entry_x.findall('ent_seq')
        entry.seqid = ent_seq_set_x[0].text
        entry.pub_date = datetime.datetime.now()
        entry.save()
        #get kanji and reading element, kanji element may be empty
        k_ele_set = entry_x.findall('k_ele') # Kanji Element
        r_ele_set = entry_x.findall('r_ele') # Reading Element ( furigana )
        #get sense elements, these include tags and definitions
        sense_x_set = entry_x.findall('sense') # Sense element
        add_keb(entry, k_ele_set)
        add_reb(entry, r_ele_set)
        add_sense(entry, sense_x_set)
        
    return list_of_entries # return the list of uploaded entries ( KEB/REB )

#Add all Kunyomi elements. 
#Do no duplicate onyomi, so we can easily find kanjis with the same kunyomi readings
def add_kunyomi(entry, kunyomi_set):
    for kunyomi_x in kunyomi_set:
        kunyomis = Kunyomi.objects.all().filter(text = kunyomi_x.text) #for entry to be duplicate the text must be exactly the same
        if not kunyomis:
            logger.info("No duplicates detected, Creating a new Onyomi")
            kunyomi = Kunyomi()
            kunyomi.text = kunyomi_x.text
            kunyomi.save()
            logger.info('ADDED NEW Onyomi: %s' % kunyomi.text)
        else :
            kunyomi = kunyomis[0]
            logger.info('WARNING:  Detected a Duplicate Onyomi in a KanjiEntry: %s. Ignoring this Onyomi' % kunyomi.text)

        kunyomi.entry.add(entry) # add entry pointer to list

#Add all Onyomi elements. 
#Do no duplicate onyomi, so we can easily find kanjis with the same onyomi readings
def add_onyomi(entry, onyomi_set):
    for onyomi_x in onyomi_set:
        onyomis = Onyomi.objects.all().filter(text = onyomi_x.text) #for entry to be duplicate the text must be exactly the same
        if not onyomis:
            logger.info("No duplicates detected, Creating a new Onyomi")
            onyomi = Onyomi()
            onyomi.text = onyomi_x.text
            onyomi.save()
            logger.info('ADDED NEW Onyomi: %s' % onyomi.text)
        else :
            onyomi = onyomis[0]
            logger.info('WARNING:  Detected a Duplicate Onyomi in a KanjiEntry: %s. Ignoring this Onyomi' % onyomi.text)

        onyomi.entry.add(entry) # add entry pointer to list

#Add all KanjiDescription elements. 
#Do no duplicate descriptions, so we can easily find kanjis with the same meaning later on
def add_kanji_description(entry, description_set):
    for description_x in description_set:
        descriptions = KanjiDescription.objects.all().filter(text = description_x.text) #for entry to be duplicate the text must be exactly the same
        if not descriptions:
            logger.info("No duplicates detected, Creating a new KanjiDescription")
            description = KanjiDescription()
            description.text = description_x.text
            description.save()
            logger.info('ADDED NEW KanjiDescription: %s' % description.text)
        else :
            description = descriptions[0]
            logger.info('WARNING:  Detected a Duplicate KanjiDescription in a KanjiEntry: %s. Ignoring this FormulaEntry' % description.text)

        #check if kanji entry already has this entry pointer?
        description.entry.add(entry) # add entry pointer to list

#Process a GKgrammar file
def process_kanji_xml_file(xml_file):
    list_of_entries = []
    number = 0

    logger2.info('Processing XML File - ')
    #Main execution code
    #load the entire GKgrammar XML File 
    tree = metree.parse(xml_file)

    #get root of the XML File
    root = tree.getroot()

    #TODO: Check the XML format...
    logger2.info('TODO: Check XML format  - ')

    #get all entries in the file
    entries_x = root.findall('kanji_entry')
    logger.info("Got root %s" % root)
    for entry_x in entries_x:
        number += 1
        #create the entry
        level_x = 0 #default level 0, level assigned in later processing
        #Check if the KanjiEntry already exists
        ent_seq_set_x = entry_x.findall('ent_seq')
        kanjientries = KanjiEntry.objects.all().filter(seqid = ent_seq_set_x[0].text) #for entry to be duplicate the text must be exactly the same
        if len(kanjientries) > 1 :
            logger.info("ERROR: Found multiple KanjiEntries with same seqid in the database ( ent_seq) %s" % ent_seq_set_x[0].text)
            entry = kanjientries[0] #default to first one found
        elif  len(kanjientries ) == 1 :
            entry = kanjientries[0] #load the existing KanjiEntry to check for updates
            logger.info("UPDATING KANJI ENTRY: Found existing KanjiEntry, with same seqid ( ent_seq) %s" % ent_seq_set_x[0].text)
        else : #create a new entry..
            entry = KanjiEntry()
            entry.seqid = ent_seq_set_x[0].text
        #blindly update the entry
        name_set_x = entry_x.findall('name')
        entry.name = name_set_x[0].text
        kanji_set_x = entry_x.findall('kanji')
        entry.text = name_set_x[0].text
        jlpt_set_x = entry_x.findall('jlpt')
        jlpt_level_set = entry_x.findall('jlpt') # Kanji Element JLPT Level
        entry.jlpt = jlpt_level_set[0].text
        level_set = entry_x.findall('level') # Sublevel
        entry.level = jlpt_level_set[0].text
        entry.save()

        description_set = entry_x.findall('description') # Full Description
        onyomi_set = entry_x.findall('onyomi') # Full Description
        kunyomi_set = entry_x.findall('kunyomi') # Full Description

        if len(description_set) >= 1 : 
            add_kanji_description(entry, description_set)
        if len(onyomi_set) >= 1 : 
            add_onyomi(entry, onyomi_set)
        if len(kunyomi_set) >= 1 : 
            add_kunyomi(entry, kunyomi_set)
        
    return list_of_entries # return the list of uploaded entries ( KEB/REB )
